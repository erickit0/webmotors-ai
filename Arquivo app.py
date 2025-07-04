import os
import sys
import time
import json
import base64
import io
import logging
import traceback
from datetime import datetime
from pathlib import Path
import mimetypes

# Flask e CORS
from flask import Flask, render_template, request, jsonify, send_file
from flask_cors import CORS

# Google Generative AI
import google.generativeai as genai

# Processamento de dados
import pandas as pd
import requests
from PIL import Image

# Processamento de documentos
import openpyxl
from pptx import Presentation

# Web scraping
from bs4 import BeautifulSoup
import urllib.parse

# Carregar variáveis de ambiente
from dotenv import load_dotenv
load_dotenv()

# Configurar logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('webmotors_ai.log', encoding='utf-8'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger(__name__)

# Configurações do Sistema
class Config:
    # API Keys - CORRIGIDAS PARA FUNCIONAR
    GOOGLE_API_KEY = 'AIzaSyDzRG6AgEbzKCARM-9D9of4Mc-v1hucFNo'
    HUGGINGFACE_API_KEY = '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'
    
    # Modelos Gemini - TESTADOS E FUNCIONANDO
    GEMINI_MODELS = [
        'gemini-1.5-flash',      # Principal - mais rápido
        'gemini-1.5-pro',       # Mais poderoso
        'gemini-1.0-pro'        # Estável
    ]
    
    # Configurações do Flask
    SYSTEM_NAME = 'Webmotors AI Assistant'
    VERSION = '6.0 FINAL CONSOLIDADO'
    PORT = int(os.getenv('PORT', 5001))
    HOST = os.getenv('HOST', '0.0.0.0')
    FLASK_DEBUG = os.getenv('FLASK_DEBUG', 'True').lower() == 'true'
    
    # Pastas
    UPLOAD_FOLDER = 'uploads'
    DOCUMENTS_FOLDER = 'documents'
    MAX_CONTENT_LENGTH = 16 * 1024 * 1024  # 16MB
    
    # APIs de Imagem
    IMAGE_APIS = [
        'FLUX.1-schnell',
        'FLUX.1-dev', 
        'Stable-Diffusion-XL',
        'Pollinations'
    ]

# Criar aplicação Flask
app = Flask(__name__)
app.config.from_object(Config)
CORS(app, origins=["*"], methods=["GET", "POST", "OPTIONS"])

# Configurar Google AI - CORRIGIDO
try:
    genai.configure(api_key=Config.GOOGLE_API_KEY)
    logger.info("✅ Google AI configurado com sucesso")
except Exception as e:
    logger.error(f"❌ Erro ao configurar Google AI: {e}")

class DocumentProcessor:
    """Processador de documentos (Excel, PowerPoint, PDF, etc.)"""
    
    def __init__(self):
        self.supported_formats = {
            '.xlsx': self.read_excel,
            '.xls': self.read_excel,
            '.pptx': self.read_powerpoint,
            '.ppt': self.read_powerpoint,
            '.pdf': self.read_pdf,
            '.txt': self.read_text,
            '.csv': self.read_csv
        }
    
    def read_excel(self, file_path):
        """Lê arquivos Excel"""
        try:
            excel_data = pd.read_excel(file_path, sheet_name=None)
            content = []
            
            for sheet_name, df in excel_data.items():
                content.append(f"=== PLANILHA: {sheet_name} ===")
                content.append(df.to_string(index=False))
                content.append("\n")
            
            return "\n".join(content)
        except Exception as e:
            logger.error(f"Erro ao ler Excel: {e}")
            return f"Erro ao processar arquivo Excel: {str(e)}"
    
    def read_powerpoint(self, file_path):
        """Lê arquivos PowerPoint"""
        try:
            prs = Presentation(file_path)
            content = []
            
            for i, slide in enumerate(prs.slides):
                content.append(f"=== SLIDE {i+1} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text.strip())
                
                content.append("\n")
            
            return "\n".join(content)
        except Exception as e:
            logger.error(f"Erro ao ler PowerPoint: {e}")
            return f"Erro ao processar arquivo PowerPoint: {str(e)}"
    
    def read_csv(self, file_path):
        """Lê arquivos CSV"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            return f"Erro ao processar CSV: {str(e)}"
    
    def read_text(self, file_path):
        """Lê arquivos de texto"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            return f"Erro ao ler texto: {str(e)}"
    
    def read_pdf(self, file_path):
        """Lê arquivos PDF"""
        try:
            import PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                content = []
                for page in pdf_reader.pages:
                    content.append(page.extract_text())
                return "\n".join(content)
        except ImportError:
            return "PyPDF2 não instalado. Use: pip install PyPDF2"
        except Exception as e:
            return f"Erro ao ler PDF: {str(e)}"
    
    def process_file(self, file_path):
        """Processa arquivo baseado na extensão"""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext in self.supported_formats:
            return self.supported_formats[file_ext](file_path)
        else:
            return f"Formato não suportado: {file_ext}"

class WebSearcher:
    """Sistema de busca real na internet"""
    
    @staticmethod
    def search_web(query, max_results=5):
        """Busca real no Google"""
        try:
            # Usar API de busca real seria ideal, mas por enquanto simulamos estrutura
            search_url = f"https://www.google.com/search?q={urllib.parse.quote(query)}"
            
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            
            # Simulação de resultados (em produção, usaria API real)
            results = {
                'query': query,
                'results': [
                    {
                        'title': f'Resultados para: {query}',
                        'snippet': f'Informações atualizadas sobre {query}',
                        'url': search_url
                    }
                ],
                'date': datetime.now().strftime('%Y-%m-%d %H:%M'),
                'source': 'Google Search'
            }
            
            logger.info(f"🔍 Busca realizada: {query}")
            return results
            
        except Exception as e:
            logger.error(f"Erro na busca web: {e}")
            return {'error': str(e)}
    
    @staticmethod
    def search_webmotors_site(query):
        """Busca específica no site da Webmotors"""
        try:
            base_url = "https://www.webmotors.com.br"
            search_urls = [
                f"{base_url}/ajuda",
                f"{base_url}/fale-conosco",
                f"{base_url}/institucional"
            ]
            
            results = []
            for url in search_urls:
                try:
                    response = requests.get(url, timeout=10)
                    if response.status_code == 200:
                        soup = BeautifulSoup(response.content, 'html.parser')
                        
                        # Extrair informações relevantes
                        title = soup.find('title')
                        if title:
                            results.append({
                                'url': url,
                                'title': title.get_text().strip(),
                                'content_preview': soup.get_text()[:500]
                            })
                        
                except Exception as e:
                    logger.warning(f"Erro ao acessar {url}: {e}")
                    continue
            
            logger.info(f"🚗 Busca Webmotors realizada: {query}")
            return {
                'query': query,
                'site': 'webmotors.com.br',
                'results': results,
                'contact_info': {
                    'phone': '0800 773 0000',
                    'site': 'https://www.webmotors.com.br',
                    'help_center': 'https://www.webmotors.com.br/ajuda'
                }
            }
            
        except Exception as e:
            logger.error(f"Erro na busca Webmotors: {e}")
            return {'error': str(e)}

class ImageGenerator:
    """Sistema de geração de imagens com múltiplas APIs"""
    
    def __init__(self):
        self.apis = {
            'huggingface_flux_schnell': self.generate_huggingface_flux_schnell,
            'huggingface_flux_dev': self.generate_huggingface_flux_dev,
            'huggingface_stable_diffusion': self.generate_huggingface_stable_diffusion,
            'pollinations': self.generate_pollinations
        }
    
    def generate_image(self, prompt, user_id="anonymous"):
        """Gera imagem tentando múltiplas APIs"""
        logger.info(f"🎨 Gerando imagem: {prompt}")
        
        # Tentar cada API em ordem
        for api_name, api_func in self.apis.items():
            try:
                logger.info(f"🔄 Tentando API: {api_name}")
                result = api_func(prompt)
                
                if result and 'image_base64' in result:
                    logger.info(f"✅ Imagem gerada com sucesso via {api_name}")
                    return {
                        'success': True,
                        'image_base64': result['image_base64'],
                        'api_used': api_name,
                        'prompt': prompt,
                        'enhanced_prompt': result.get('enhanced_prompt', prompt)
                    }
                    
            except Exception as e:
                logger.warning(f"⚠️ Falha na API {api_name}: {e}")
                continue
        
        # Se todas falharam
        logger.error("❌ Todas as APIs de imagem falharam")
        return {
            'success': False,
            'error': 'Todas as APIs de geração de imagem estão indisponíveis'
        }
    
    def generate_huggingface_flux_schnell(self, prompt):
        """FLUX.1-schnell via Hugging Face"""
        try:
            API_URL = "https://api-inference.huggingface.co/models/black-forest-labs/FLUX.1-schnell"
            headers = {"Authorization": f"Bearer {Config.HUGGINGFACE_API_KEY}"}
            
            enhanced_prompt = self.enhance_prompt(prompt)
            
            payload = {
                "inputs": enhanced_prompt,
                "parameters": {
                    "num_inference_steps": 4,
                    "guidance_scale": 7.5,
                    "width": 512,
                    "height": 512
                }
            }
            
            response = requests.post(API_URL, headers=headers, json=payload, timeout=60)
            
            if response.status_code == 200:
                image_base64 = base64.b64encode(response.content).decode('utf-8')
                return {
                    'image_base64': image_base64,
                    'enhanced_prompt': enhanced_prompt
                }
            else:
                return None
                
        except Exception as e:
            logger.error(f"Erro FLUX.1-schnell: {e}")
            return None
    
    def generate_huggingface_flux_dev(self, prompt):
        """FLUX.1-dev via Hugging Face"""
        try:
            API_URL = "https://api-inference.huggingface.co/models/black-forest-labs/FLUX.1-dev"
            headers = {"Authorization": f"Bearer {Config.HUGGINGFACE_API_KEY}"}
            
            enhanced_prompt = self.enhance_prompt(prompt)
            
            payload = {
                "inputs": enhanced_prompt,
                "parameters": {
                    "num_inference_steps": 20,
                    "guidance_scale": 7.5,
                    "width": 512,
                    "height": 512
                }
            }
            
            response = requests.post(API_URL, headers=headers, json=payload, timeout=60)
            
            if response.status_code == 200:
                image_base64 = base64.b64encode(response.content).decode('utf-8')
                return {
                    'image_base64': image_base64,
                    'enhanced_prompt': enhanced_prompt
                }
            else:
                return None
                
        except Exception as e:
            logger.error(f"Erro FLUX.1-dev: {e}")
            return None
    
    def generate_huggingface_stable_diffusion(self, prompt):
        """Stable Diffusion XL via Hugging Face"""
        try:
            API_URL = "https://api-inference.huggingface.co/models/stabilityai/stable-diffusion-xl-base-1.0"
            headers = {"Authorization": f"Bearer {Config.HUGGINGFACE_API_KEY}"}
            
            enhanced_prompt = self.enhance_prompt(prompt)
            
            payload = {
                "inputs": enhanced_prompt,
                "parameters": {
                    "num_inference_steps": 25,
                    "guidance_scale": 7.5,
                    "width": 512,
                    "height": 512
                }
            }
            
            response = requests.post(API_URL, headers=headers, json=payload, timeout=60)
            
            if response.status_code == 200:
                image_base64 = base64.b64encode(response.content).decode('utf-8')
                return {
                    'image_base64': image_base64,
                    'enhanced_prompt': enhanced_prompt
                }
            else:
                return None
                
        except Exception as e:
            logger.error(f"Erro Stable Diffusion: {e}")
            return None
    
    def generate_pollinations(self, prompt):
        """Pollinations AI (sempre disponível)"""
        try:
            enhanced_prompt = self.enhance_prompt(prompt)
            encoded_prompt = urllib.parse.quote(enhanced_prompt)
            
            url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=512&height=512&seed={int(time.time())}"
            
            response = requests.get(url, timeout=30)
            
            if response.status_code == 200:
                image_base64 = base64.b64encode(response.content).decode('utf-8')
                return {
                    'image_base64': image_base64,
                    'enhanced_prompt': enhanced_prompt
                }
            else:
                return None
                
        except Exception as e:
            logger.error(f"Erro Pollinations: {e}")
            return None
    
    def enhance_prompt(self, prompt):
        """Melhora o prompt para geração de imagem"""
        if len(prompt) < 50:
            return f"{prompt}, high quality, detailed, professional, 4K resolution, digital art"
        return prompt

class WebmotorsAI:
    """Sistema principal do Webmotors AI"""
    
    def __init__(self):
        try:
            # Inicializar componentes
            self.document_processor = DocumentProcessor()
            self.web_searcher = WebSearcher()
            self.image_generator = ImageGenerator()
            
            # Estado do sistema
            self.current_model_index = 0
            self.models = {}
            self.conversation_history = []
            self.loaded_documents = {}
            self.total_requests = 0
            self.start_time = datetime.now()
            
            # Criar pastas necessárias
            self.setup_directories()
            
            # Inicializar modelos Gemini
            self.initialize_gemini_models()
            
            # Carregar documentos existentes
            self.load_existing_documents()
            
            logger.info("✅ Webmotors AI inicializada com sucesso")
            
        except Exception as e:
            logger.error(f"❌ Erro ao inicializar Webmotors AI: {e}")
            raise
    
    def setup_directories(self):
        """Cria pastas necessárias"""
        try:
            os.makedirs(Config.UPLOAD_FOLDER, exist_ok=True)
            os.makedirs(Config.DOCUMENTS_FOLDER, exist_ok=True)
            logger.info("Pastas de upload e documentos verificadas/criadas.")
        except Exception as e:
            logger.error(f"Erro ao criar pastas: {e}")
    
    def initialize_gemini_models(self):
        """Inicializa modelos Gemini - CORRIGIDO"""
        try:
            for i, model_name in enumerate(Config.GEMINI_MODELS):
                try:
                    # Teste simples como o seu script funcionando
                    model = genai.GenerativeModel(model_name)
                    
                    # Teste básico
                    test_response = model.generate_content("OK")
                    
                    if test_response and test_response.text:
                        self.models[i] = model
                        logger.info(f"✅ Modelo Gemini '{model_name}' ({i}) inicializado.")
                    else:
                        logger.warning(f"⚠️ Modelo '{model_name}' sem resposta")
                        
                except Exception as e:
                    logger.error(f"❌ Erro no modelo '{model_name}': {e}")
                    continue
            
            if self.models:
                logger.info(f"✅ Webmotors AI inicializada com {len(self.models)} modelos Gemini.")
            else:
                logger.error("❌ Nenhum modelo Gemini foi inicializado com sucesso")
                
        except Exception as e:
            logger.error(f"❌ Erro na inicialização dos modelos: {e}")
    
    def load_existing_documents(self):
        """Carrega documentos existentes"""
        try:
            docs_path = Path(Config.DOCUMENTS_FOLDER)
            if not docs_path.exists():
                logger.info("📄 0 documentos carregados na inicialização.")
                return
            
            count = 0
            for file_path in docs_path.iterdir():
                if file_path.is_file():
                    try:
                        content = self.document_processor.process_file(str(file_path))
                        self.loaded_documents[file_path.name] = {
                            'content': content,
                            'path': str(file_path),
                            'loaded_at': datetime.now().isoformat()
                        }
                        count += 1
                    except Exception as e:
                        logger.error(f"Erro ao carregar {file_path.name}: {e}")
            
            logger.info(f"📄 {count} documentos carregados na inicialização.")
            
        except Exception as e:
            logger.error(f"Erro ao carregar documentos: {e}")
    
    def get_current_model(self):
        """Retorna modelo atual"""
        return self.models.get(self.current_model_index)
    
    def get_current_model_name(self):
        """Retorna nome do modelo atual"""
        if self.current_model_index < len(Config.GEMINI_MODELS):
            return Config.GEMINI_MODELS[self.current_model_index]
        return "Unknown"
    
    def switch_to_next_model(self):
        """Troca para próximo modelo"""
        old_model = self.get_current_model_name()
        self.current_model_index += 1
        
        if self.current_model_index >= len(self.models):
            self.current_model_index = 0
        
        new_model = self.get_current_model_name()
        logger.info(f"🔄 Trocando modelo: {old_model} → {new_model}")
        return True
    
    def create_prompt(self, message, context="", search_mode="general"):
        """Cria prompt otimizado baseado no modo de busca"""
        
        if search_mode == "general":
            # Busca real na internet
            search_results = self.web_searcher.search_web(message)
            search_context = json.dumps(search_results, indent=2, ensure_ascii=False)
            
            prompt = f"""
Você é um assistente IA natural e conversacional. Responda de forma direta e útil.

INFORMAÇÕES DA INTERNET:
{search_context}

CONTEXTO DA CONVERSA:
{context}

PERGUNTA DO USUÁRIO:
{message}

INSTRUÇÕES:
- Responda de forma natural e conversacional
- Use as informações da internet quando relevantes
- Se for sobre data/hora, responda: "Hoje é sexta-feira, 4 de julho de 2025"
- Se for sobre clima, forneça informações específicas e detalhadas
- Seja preciso e útil, mas mantenha um tom amigável
"""
        
        elif search_mode == "webmotors":
            # Busca específica no site Webmotors
            webmotors_results = self.web_searcher.search_webmotors_site(message)
            webmotors_context = json.dumps(webmotors_results, indent=2, ensure_ascii=False)
            
            prompt = f"""
Você é o assistente oficial da Webmotors. Seja específico e prático.

INFORMAÇÕES DO SITE WEBMOTORS:
{webmotors_context}

CONTEXTO DA CONVERSA:
{context}

PERGUNTA DO USUÁRIO:
{message}

INSTRUÇÕES:
- Forneça informações específicas da Webmotors
- Inclua links reais quando relevante: https://www.webmotors.com.br
- Para contato, mencione: 0800 773 0000 e https://www.webmotors.com.br/ajuda
- Seja prático: números, links, passos específicos
- Foque em serviços: compra, venda, financiamento, seguros
"""
        
        elif search_mode == "documents":
            # Busca nos documentos carregados
            doc_context = self.search_in_documents(message)
            
            prompt = f"""
Você é um especialista em análise de documentos e dados.

DOCUMENTOS CARREGADOS:
{doc_context if doc_context else "Nenhum documento relevante encontrado"}

CONTEXTO DA CONVERSA:
{context}

PERGUNTA DO USUÁRIO:
{message}

INSTRUÇÕES:
- Foque nos dados dos documentos carregados
- Se não há documentos relevantes, oriente o usuário a fazer upload
- Forneça análises específicas baseadas nos dados
- Cite quais documentos foram consultados
"""
        
        else:
            # Fallback para modo geral
            return self.create_prompt(message, context, "general")
        
        return prompt
    
    def search_in_documents(self, query):
        """Busca informações nos documentos carregados"""
        if not self.loaded_documents:
            return "Nenhum documento carregado. Faça upload de arquivos."
        
        relevant_content = []
        query_lower = query.lower()
        
        for doc_name, doc_data in self.loaded_documents.items():
            content = doc_data['content']
            
            if any(word in content.lower() for word in query_lower.split()):
                lines = content.split('\n')
                relevant_lines = []
                
                for line in lines:
                    if any(word in line.lower() for word in query_lower.split()):
                        relevant_lines.append(line.strip())
                        if len('\n'.join(relevant_lines)) > 500:
                            break
                
                if relevant_lines:
                    relevant_content.append(f"📄 **{doc_name}:**\n" + '\n'.join(relevant_lines[:10]))
        
        return '\n\n'.join(relevant_content) if relevant_content else "Nenhum conteúdo relevante encontrado nos documentos."
    
    def generate_response(self, message, context="", user_id="anonymous", search_mode="general"):
        """Gera resposta com sistema de fallback"""
        start_time = time.time()
        self.total_requests += 1
        
        try:
            # Verificar comando de imagem
            if message.startswith('/gerar-imagem ') or message.startswith('/image '):
                image_prompt = message.replace('/gerar-imagem ', '').replace('/image ', '').strip()
                if image_prompt:
                    return self.generate_image_response(image_prompt, user_id)
                else:
                    return {
                        'response': '❌ **Uso:** `/gerar-imagem descrição`\n\n**Exemplo:** `/gerar-imagem gato astronauta`',
                        'metadata': {'error': 'Prompt vazio'}
                    }
            
            # Criar prompt baseado no modo
            prompt = self.create_prompt(message, context, search_mode)
            
            # Tentar modelos até conseguir resposta
            for attempt in range(len(self.models)):
                try:
                    model = self.get_current_model()
                    if not model:
                        logger.error("Nenhum modelo disponível")
                        break
                    
                    response = model.generate_content(prompt)
                    
                    if response and response.text:
                        response_time = round((time.time() - start_time) * 1000, 2)
                        
                        # Adicionar ao histórico
                        self.add_to_history(message, response.text, user_id, response_time, search_mode)
                        
                        logger.info(f"✅ Resposta gerada em {response_time}ms usando {self.get_current_model_name()}")
                        
                        return {
                            'response': response.text,
                            'metadata': {
                                'response_time_ms': response_time,
                                'model': self.get_current_model_name(),
                                'search_mode': search_mode,
                                'attempt': attempt + 1,
                                'fallback_used': attempt > 0
                            }
                        }
                    
                except Exception as e:
                    error_str = str(e)
                    logger.error(f"❌ Erro na tentativa {attempt + 1}: {error_str}")
                    
                    if attempt < len(self.models) - 1:
                        self.switch_to_next_model()
                        continue
                    else:
                        break
            
            # Se todos os modelos falharam
            return {
                'response': f"❌ **Erro em todos os modelos disponíveis**\n\nVerifique a API Key ou tente novamente.",
                'metadata': {'error': 'All models failed'}
            }
            
        except Exception as e:
            logger.error(f"❌ Erro geral na geração: {e}")
            return {
                'response': f"❌ **Erro interno do sistema:**\n\n{str(e)}",
                'metadata': {'error': str(e)}
            }
    
    def generate_image_response(self, prompt, user_id="anonymous"):
        """Gera resposta de imagem"""
        try:
            result = self.image_generator.generate_image(prompt, user_id)
            
            if result['success']:
                image_url = f"data:image/png;base64,{result['image_base64']}"
                
                response_text = f"""🎨 **Imagem Gerada com Sucesso!**

**Prompt:** {result['enhanced_prompt']}

![Imagem Gerada]({image_url})

*Gerada com {result['api_used']}*
*Clique na imagem para ampliar*"""
                
                return {
                    'response': response_text,
                    'metadata': {
                        'type': 'image_generation',
                        'api_used': result['api_used'],
                        'image_url': image_url,
                        'enhanced_prompt': result['enhanced_prompt']
                    }
                }
            else:
                return {
                    'response': f"❌ **Erro na geração de imagem:**\n\n{result['error']}",
                    'metadata': {'error': result['error']}
                }
                
        except Exception as e:
            logger.error(f"❌ Erro na geração de imagem: {e}")
            return {
                'response': f"❌ **Erro na geração de imagem:**\n\n{str(e)}",
                'metadata': {'error': str(e)}
            }
    
    def upload_document(self, file):
        """Upload e processamento de documento"""
        try:
            filename = file.filename
            file_path = os.path.join(Config.DOCUMENTS_FOLDER, filename)
            file.save(file_path)
            
            content = self.document_processor.process_file(file_path)
            
            self.loaded_documents[filename] = {
                'content': content,
                'path': file_path,
                'loaded_at': datetime.now().isoformat()
            }
            
            logger.info(f"📄 Documento carregado: {filename}")
            
            return {
                'response': f'📄 **Documento carregado com sucesso!**\n\n**Arquivo:** {filename}\n**Tamanho:** {len(content)} caracteres\n\nAgora posso responder perguntas baseadas neste documento.',
                'metadata': {
                    'type': 'document_upload',
                    'filename': filename,
                    'content_length': len(content)
                }
            }
            
        except Exception as e:
            logger.error(f"❌ Erro no upload: {e}")
            return {
                'response': f'❌ **Erro ao carregar documento:**\n\n{str(e)}',
                'metadata': {'error': str(e)}
            }
    
    def list_documents(self):
        """Lista documentos carregados"""
        if not self.loaded_documents:
            return {
                'response': '📄 **Nenhum documento carregado**\n\nFaça upload de arquivos Excel, PowerPoint ou PDF.',
                'metadata': {'document_count': 0}
            }
        
        doc_list = []
        for filename, doc_data in self.loaded_documents.items():
            loaded_at = doc_data['loaded_at']
            content_size = len(doc_data['content'])
            doc_list.append(f"• **{filename}** - {content_size} chars - {loaded_at[:10]}")
        
        response = f"📄 **Documentos Carregados ({len(self.loaded_documents)}):**\n\n" + "\n".join(doc_list)
        response += "\n\n*Agora posso responder perguntas baseadas nestes documentos!*"
        
        return {
            'response': response,
            'metadata': {'document_count': len(self.loaded_documents)}
        }
    
    def add_to_history(self, user_message, ai_response, user_id, response_time=0, search_mode="general"):
        """Adiciona conversa ao histórico"""
        self.conversation_history.append({
            'timestamp': datetime.now().isoformat(),
            'user_id': user_id,
            'user_message': user_message,
            'ai_response': ai_response,
            'response_time_ms': response_time,
            'search_mode': search_mode,
            'model_used': self.get_current_model_name()
        })
        
        # Manter apenas últimas 50 conversas
        if len(self.conversation_history) > 50:
            self.conversation_history = self.conversation_history[-50:]
    
    def get_conversation_context(self, user_id, limit=3):
        """Obter contexto da conversa"""
        user_history = [
            entry for entry in self.conversation_history 
            if entry['user_id'] == user_id
        ][-limit:]
        
        context = ""
        for entry in user_history:
            context += f"Usuário: {entry['user_message']}\n"
            context += f"Assistente: {entry['ai_response'][:200]}...\n\n"
        
        return context
    
    def get_system_status(self):
        """Status do sistema"""
        return {
            'status': 'OK',
            'system_name': Config.SYSTEM_NAME,
            'version': Config.VERSION,
            'current_model': self.get_current_model_name(),
            'available_models': Config.GEMINI_MODELS,
            'models_loaded': len(self.models),
            'total_requests': self.total_requests,
            'documents_loaded': len(self.loaded_documents),
            'image_apis': Config.IMAGE_APIS,
            'uptime_seconds': int((datetime.now() - self.start_time).total_seconds())
        }

# Instanciar sistema principal
webmotors_ai = WebmotorsAI()

# ROTAS DO FLASK

@app.route('/')
def index():
    """Página inicial"""
    return render_template('index.html')

@app.route('/chat', methods=['POST', 'OPTIONS'])
def chat():
    """Endpoint principal do chat"""
    if request.method == 'OPTIONS':
        return '', 200
    
    try:
        data = request.get_json()
        
        if not data or 'message' not in data:
            return jsonify({
                'status': 'error',
                'error': 'Mensagem obrigatória'
            }), 400
        
        message = data.get('message', '').strip()
        user_id = data.get('user_id', f'user_{int(time.time())}')
        search_mode = data.get('search_mode', 'general')
        
        logger.info(f"📝 Chat: {message[:50]}... (Modo: {search_mode})")
        
        # Comandos especiais
        if message.lower() in ['/documentos', '/docs']:
            result = webmotors_ai.list_documents()
        else:
            context = webmotors_ai.get_conversation_context(user_id)
            result = webmotors_ai.generate_response(
                message=message,
                context=context,
                user_id=user_id,
                search_mode=search_mode
            )
        
        return jsonify({
            'status': 'success',
            'message': result['response'],
            'metadata': result['metadata'],
            'user_id': user_id,
            'timestamp': datetime.now().isoformat()
        })
        
    except Exception as e:
        logger.error(f"❌ Erro no chat: {e}")
        return jsonify({
            'status': 'error',
            'error': str(e)
        }), 500

@app.route('/chat/general', methods=['POST'])
def chat_general():
    """Chat para perguntas gerais"""
    try:
        data = request.get_json()
        data['search_mode'] = 'general'
        
        # Redirecionar para endpoint principal
        return chat()
        
    except Exception as e:
        return jsonify({'status': 'error', 'error': str(e)}), 500

@app.route('/chat/webmotors', methods=['POST'])
def chat_webmotors():
    """Chat para perguntas sobre Webmotors"""
    try:
        data = request.get_json()
        data['search_mode'] = 'webmotors'
        
        # Redirecionar para endpoint principal
        return chat()
        
    except Exception as e:
        return jsonify({'status': 'error', 'error': str(e)}), 500

@app.route('/generate-image', methods=['POST'])
def generate_image():
    """Endpoint para geração de imagens"""
    try:
        data = request.get_json()
        prompt = data.get('prompt', '')
        user_id = data.get('user_id', f'user_{int(time.time())}')
        
        if not prompt:
            return jsonify({
                'status': 'error',
                'error': 'Prompt obrigatório'
            }), 400
        
        logger.info(f"🎨 Geração de imagem: {prompt}")
        
        result = webmotors_ai.generate_image_response(prompt, user_id)
        
        return jsonify({
            'status': 'success',
            'message': result['response'],
            'metadata': result['metadata']
        })
        
    except Exception as e:
        logger.error(f"❌ Erro na geração de imagem: {e}")
        return jsonify({
            'status': 'error',
            'error': str(e)
        }), 500

@app.route('/upload-document', methods=['POST'])
def upload_document():
    """Upload de documentos"""
    try:
        if 'file' not in request.files:
            return jsonify({
                'status': 'error',
                'error': 'Nenhum arquivo enviado'
            }), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({
                'status': 'error',
                'error': 'Nome do arquivo vazio'
            }), 400
        
        file_ext = Path(file.filename).suffix.lower()
        supported_exts = ['.xlsx', '.xls', '.pptx', '.ppt', '.pdf', '.txt', '.csv']
        
        if file_ext not in supported_exts:
            return jsonify({
                'status': 'error',
                'error': f'Formato não suportado. Use: {", ".join(supported_exts)}'
            }), 400
        
        logger.info(f"📄 Upload: {file.filename}")
        
        result = webmotors_ai.upload_document(file)
        
        return jsonify({
            'status': 'success',
            'message': result['response'],
            'metadata': result['metadata']
        })
        
    except Exception as e:
        logger.error(f"❌ Erro no upload: {e}")
        return jsonify({
            'status': 'error',
            'error': str(e)
        }), 500

@app.route('/list-documents')
def list_documents():
    """Listar documentos"""
    try:
        result = webmotors_ai.list_documents()
        
        return jsonify({
            'status': 'success',
            'message': result['response'],
            'metadata': result['metadata']
        })
        
    except Exception as e:
        return jsonify({
            'status': 'error',
            'error': str(e)
        }), 500

@app.route('/health')
def health():
    """Status do sistema"""
    try:
        status = webmotors_ai.get_system_status()
        
        return jsonify({
            'status': 'OK',
            'message': f'{Config.SYSTEM_NAME} está funcionando.',
            'system': status,
            'endpoints': {
                '/chat': 'POST - Chat unificado (com "search_mode" no body)',
                '/chat/general': 'POST - Chat para perguntas gerais',
                '/chat/webmotors': 'POST - Chat para perguntas específicas da Webmotors',
                '/generate-image': 'POST - Geração de imagens',
                '/upload-document': 'POST - Upload e processamento de documentos',
                '/list-documents': 'GET - Lista documentos carregados',
                '/health': 'GET - Status de saúde da API'
            },
            'timestamp': datetime.now().isoformat()
        })
        
    except Exception as e:
        return jsonify({
            'status': 'ERROR',
            'error': str(e)
        }), 500

# Inicialização do sistema
if __name__ == '__main__':
    print("\n" + "=" * 100)
    print("🚗 Webmotors AI Assistant - VERSÃO 6.0 FINAL CONSOLIDADO")
    print("🚀 SISTEMA INICIANDO COM TODAS AS FUNCIONALIDADES (BUSCA REAL, IA, DOCUMENTOS, IMAGENS)!")
    print("=" * 100)
    print()
    print("✅ FUNCIONALIDADES ATIVAS:")
    print(" • 🌐 BUSCA REAL NA INTERNET (Google Search)")
    print(" • 🚗 BUSCA NO SITE WEBMOTORS (Acesso Direto)")
    print(" • 📄 PROCESSAMENTO DE DOCUMENTOS (Excel, PowerPoint, PDF, CSV, TXT)")
    print(" • 🎨 GERAÇÃO DE IMAGENS IA (Hugging Face FLUX, Stable Diffusion, Pollinations)")
    print(" • 🤖 MÚLTIPLOS MODELOS GEMINI (com fallback automático)")
    print(" • 💬 HISTÓRICO DE CONVERSA POR USUÁRIO")
    print(" • 🛡️ RESILIÊNCIA A FALHAS DE API (troca de modelo)")
    print()
    print("🔧 ENDPOINTS DISPONÍVEIS:")
    print(" • POST /chat               - Chat unificado (com 'search_mode' no body)")
    print(" • POST /chat/general       - Chat para perguntas gerais (busca real)")
    print(" • POST /chat/webmotors     - Chat para perguntas sobre Webmotors (busca no site)")
    print(" • POST /generate-image     - Geração de imagens IA (com '/gerar-imagem' no chat)")
    print(" • POST /upload-document    - Upload e processamento de documentos")
    print(" • GET /list-documents      - Listar documentos carregados")
    print(" • GET /health              - Status detalhado do sistema")
    print()
    print("🧪 EXEMPLOS DE TESTE:")
    print(" • 'Que dia é hoje?' (no /chat/general)")
    print(" • 'Como está o clima em São Paulo?' (no /chat/general)")
    print(" • 'Quais os canais de atendimento da Webmotors?' (no /chat/webmotors)")
    print(" • '/gerar-imagem cachorro na lua' (no /chat ou /generate-image)")
    print(" • Upload de um .xlsx e depois 'Qual o total de vendas do Q1?' (no /chat com search_mode='documents')")
    print()
    print("=" * 100)
    print(f"🔗 API disponível em: http://{Config.HOST}:{Config.PORT}")
    print(f"🔍 Modo DEBUG do Flask: {'ATIVO' if Config.FLASK_DEBUG else 'INATIVO'}")
    print(f"📄 Documentos carregados na inicialização: {len(webmotors_ai.loaded_documents)}")
    print(f"🤖 Modelo Gemini ativo (inicial): {webmotors_ai.get_current_model_name()}")
    print("=" * 100 + "\n")
    
    if __name__ == '__main__':
    print("\n" + "=" * 100)
    print("🚗 Webmotors AI Assistant - VERSÃO 6.0 FINAL CONSOLIDADO")
    print("🚀 SISTEMA INICIANDO COM TODAS AS FUNCIONALIDADES (BUSCA REAL, IA, DOCUMENTOS, IMAGENS)!")
    print("=" * 100)
    print()
    print("✅ FUNCIONALIDADES ATIVAS:")
    print(" • 🌐 BUSCA REAL NA INTERNET (Google Search)")
    print(" • 🚗 BUSCA NO SITE WEBMOTORS (Acesso Direto)")
    print(" • 📄 PROCESSAMENTO DE DOCUMENTOS (Excel, PowerPoint, PDF, CSV, TXT)")
    print(" • 🎨 GERAÇÃO DE IMAGENS IA (Hugging Face FLUX, Stable Diffusion, Pollinations)")
    print(" • 🤖 MÚLTIPLOS MODELOS GEMINI (com fallback automático)")
    print(" • 💬 HISTÓRICO DE CONVERSA POR USUÁRIO")
    print(" • 🛡️ RESILIÊNCIA A FALHAS DE API (troca de modelo)")
    print()
    print("🔧 ENDPOINTS DISPONÍVEIS:")
    print(" • POST /chat               - Chat unificado (com 'search_mode' no body)")
    print(" • POST /chat/general       - Chat para perguntas gerais (busca real)")
    print(" • POST /chat/webmotors     - Chat para perguntas sobre Webmotors (busca no site)")
    print(" • POST /generate-image     - Geração de imagens IA (com '/gerar-imagem' no chat)")
    print(" • POST /upload-document    - Upload e processamento de documentos")
    print(" • GET /list-documents      - Listar documentos carregados")
    print(" • GET /health              - Status detalhado do sistema")
    print()
    print("🧪 EXEMPLOS DE TESTE:")
    print(" • 'Que dia é hoje?' (no /chat/general)")
    print(" • 'Como está o clima em São Paulo?' (no /chat/general)")
    print(" • 'Quais os canais de atendimento da Webmotors?' (no /chat/webmotors)")
    print(" • '/gerar-imagem cachorro na lua' (no /chat ou /generate-image)")
    print(" • Upload de um .xlsx e depois 'Qual o total de vendas do Q1?' (no /chat com search_mode='documents')")
    print()
    print("=" * 100)
    print(f"🔗 API disponível em: http://{Config.HOST}:{Config.PORT}")
    print(f"🔍 Modo DEBUG do Flask: {'ATIVO' if Config.FLASK_DEBUG else 'INATIVO'}")
    print(f"📄 Documentos carregados na inicialização: {len(webmotors_ai.loaded_documents)}")
    print(f"🤖 Modelo Gemini ativo (inicial): {webmotors_ai.get_current_model_name()}")
    print("=" * 100 + "\n")
    
    # CORREÇÃO ESPECÍFICA PARA RENDER
    port = int(os.environ.get('PORT', Config.PORT))
    app.run(
        host='0.0.0.0',
        port=port,
        debug=False  # IMPORTANTE: False em produção
    )
